from __future__ import annotations

import base64
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from autobloggy.artifacts import extract_frontmatter, format_markdown_with_frontmatter, read_json, read_yaml
from tests.helpers import copy_repo, parse_kv, resolve_generated_outline, resolve_generated_strategy, run_cli, run_cli_failure


PNG_BYTES = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO6Z0V0AAAAASUVORK5CYII="
)


def create_test_pptx(path: Path) -> None:
    image_path = path.parent / "slide-image.png"
    image_path.write_bytes(PNG_BYTES)

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1.5))
    title_box.text_frame.text = "Architecture Overview"
    body_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(1.5))
    body_box.text_frame.text = "Operators need a clear separation between raw and prepared inputs."
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(3), width=Inches(2), height=Inches(2))
    presentation.save(path)


def build_post_with_draft(repo: Path, slug: str) -> None:
    run_cli(repo, "new-post", "--slug", slug, "--topic", "Visual prep workflow")
    strategy_path = repo / "posts" / slug / "strategy.md"
    resolve_generated_strategy(strategy_path)
    run_cli(repo, "approve-strategy", "--slug", slug)
    run_cli(repo, "decide-discovery", "--slug", slug, "--decision", "no")
    run_cli(repo, "generate-outline", "--slug", slug)
    resolve_generated_outline(repo / "posts" / slug / "outline.md")
    run_cli(repo, "approve-outline", "--slug", slug)
    run_cli(repo, "generate-draft", "--slug", slug)


def test_new_post_supports_topic_only(repo_root: Path, tmp_path: Path, monkeypatch) -> None:
    repo = copy_repo(repo_root, tmp_path)
    monkeypatch.chdir(repo)

    result = parse_kv(run_cli(repo, "new-post", "--topic", "Why AI eval loops fail").stdout)
    slug = result["slug"]

    brief_path = repo / "posts" / slug / "inputs" / "user_provided" / "brief.md"
    prepared_input = repo / "posts" / slug / "inputs" / "prepared" / "input.md"
    manifest_path = repo / "posts" / slug / "inputs" / "prepared" / "input_manifest.yaml"
    strategy_path = repo / "posts" / slug / "strategy.md"
    assert slug == "why-ai-eval-loops-fail"
    assert brief_path.exists()
    assert prepared_input.exists()
    assert manifest_path.exists()
    assert strategy_path.exists()

    frontmatter, _ = extract_frontmatter(strategy_path.read_text(encoding="utf-8"))
    assert frontmatter["input_path"] == str(prepared_input)
    assert frontmatter["input_root"] == str(prepared_input.parent)
    assert frontmatter["input_manifest"] == f"posts/{slug}/inputs/prepared/input_manifest.yaml"
    assert frontmatter["preset"] == "default"


def test_new_post_with_slug_only_scaffolds_layout(repo_root: Path, tmp_path: Path, monkeypatch) -> None:
    repo = copy_repo(repo_root, tmp_path)
    monkeypatch.chdir(repo)

    slug = "scaffold-only"
    result = parse_kv(run_cli(repo, "new-post", "--slug", slug).stdout)

    post_root = repo / "posts" / slug
    assert result["status"] == "scaffolded"
    assert (post_root / "inputs" / "user_provided" / "README.md").exists()
    assert (post_root / "inputs" / "user_provided" / "brief.md").exists()
    assert (post_root / "inputs" / "user_provided" / "raw").is_dir()
    assert (post_root / "inputs" / "extracted" / "text").is_dir()
    assert (post_root / "inputs" / "prepared").is_dir()
    assert not (post_root / "strategy.md").exists()
    assert "prepare-inputs" in result["next_step"]


def test_new_post_uses_existing_post_input_folder_by_default(repo_root: Path, tmp_path: Path, monkeypatch) -> None:
    repo = copy_repo(repo_root, tmp_path)
    monkeypatch.chdir(repo)

    slug = "folder-first"
    input_root = repo / "posts" / slug / "inputs" / "user_provided"
    input_root.mkdir(parents=True, exist_ok=True)
    seeded = input_root / "notes.md"
    seeded.write_text(
        """---
title: Folder First
---

# Folder First

Operators want one place to drop source files before the writing flow starts.
""",
        encoding="utf-8",
    )

    result = parse_kv(run_cli(repo, "new-post", "--slug", slug).stdout)
    prepared_input = repo / "posts" / slug / "inputs" / "prepared" / "input.md"
    manifest_path = repo / "posts" / slug / "inputs" / "prepared" / "input_manifest.yaml"
    assert result["slug"] == slug
    assert prepared_input.exists()
    assert (repo / "posts" / slug / "strategy.md").exists()
    assert "inputs/user_provided/notes.md" in prepared_input.read_text(encoding="utf-8")
    manifest = read_yaml(manifest_path)
    assert manifest["raw_text_sources"][0]["path"] == "inputs/user_provided/notes.md"


def test_new_post_copies_source_files_into_user_raw(repo_root: Path, tmp_path: Path, monkeypatch) -> None:
    repo = copy_repo(repo_root, tmp_path)
    monkeypatch.chdir(repo)

    source_md = tmp_path / "source-notes.md"
    source_md.write_text(
        """---
title: Source Notes
---

# Source Notes

This note should be copied into the post input bundle.
""",
        encoding="utf-8",
    )
    source_txt = tmp_path / "overview.txt"
    source_txt.write_text("A second supporting file.", encoding="utf-8")

    slug = "copied-sources"
    run_cli(
        repo,
        "new-post",
        "--slug",
        slug,
        "--topic",
        "Copied sources should land in the default post input folder.",
        "--source",
        str(source_md),
        "--source",
        str(source_txt),
    )

    raw_root = repo / "posts" / slug / "inputs" / "user_provided" / "raw"
    assert (raw_root / "source-notes.md").exists()
    assert (raw_root / "overview.txt").exists()
    assert not (repo / "posts" / slug / "inputs" / "user_provided" / "supporting").exists()


def test_prepare_inputs_extracts_text_and_visuals_from_pptx(repo_root: Path, tmp_path: Path, monkeypatch) -> None:
    repo = copy_repo(repo_root, tmp_path)
    monkeypatch.chdir(repo)

    slug = "pptx-extract"
    run_cli(repo, "new-post", "--slug", slug)

    pptx_path = tmp_path / "architecture-deck.pptx"
    create_test_pptx(pptx_path)
    raw_root = repo / "posts" / slug / "inputs" / "user_provided" / "raw"
    shutil.copy2(pptx_path, raw_root / pptx_path.name)

    generated = parse_kv(run_cli(repo, "prepare-inputs", "--slug", slug).stdout)
    prepared_input = repo / "posts" / slug / "inputs" / "prepared" / "input.md"
    manifest_path = repo / "posts" / slug / "inputs" / "prepared" / "input_manifest.yaml"

    assert generated["input"] == str(prepared_input)
    assert manifest_path.exists()
    extracted_text_files = sorted((repo / "posts" / slug / "inputs" / "extracted" / "text").glob("*.md"))
    extracted_visual_files = sorted((repo / "posts" / slug / "inputs" / "extracted" / "visuals").rglob("*.png"))
    assert extracted_text_files
    assert extracted_visual_files

    manifest = read_yaml(manifest_path)
    assert manifest["extracted_text_sources"]
    assert manifest["extracted_visual_sources"]
    assert manifest["canonical_input"] == "inputs/prepared/input.md"


def test_prepare_inputs_normalizes_legacy_supporting_directory(repo_root: Path, tmp_path: Path, monkeypatch) -> None:
    repo = copy_repo(repo_root, tmp_path)
    monkeypatch.chdir(repo)

    slug = "legacy-supporting"
    post_root = repo / "posts" / slug / "inputs" / "user_provided" / "supporting"
    post_root.mkdir(parents=True, exist_ok=True)
    legacy_file = post_root / "legacy.md"
    legacy_file.write_text("# Legacy\n\nLegacy supporting material.\n", encoding="utf-8")

    run_cli(repo, "prepare-inputs", "--slug", slug)
    manifest = read_yaml(repo / "posts" / slug / "inputs" / "prepared" / "input_manifest.yaml")
    assert manifest["raw_text_sources"][0]["path"] == "inputs/user_provided/supporting/legacy.md"


def test_generate_outline_and_generate_draft_require_approvals(repo_root: Path, tmp_path: Path, monkeypatch) -> None:
    repo = copy_repo(repo_root, tmp_path)
    monkeypatch.chdir(repo)

    slug = "review-gated"
    run_cli(repo, "new-post", "--slug", slug, "--topic", "Review-gated post")
    strategy_path = repo / "posts" / slug / "strategy.md"

    blocked_outline = run_cli_failure(repo, "generate-outline", "--slug", slug)
    assert blocked_outline.returncode != 0
    assert "approve-strategy" in (blocked_outline.stdout + blocked_outline.stderr)

    resolve_generated_strategy(strategy_path)
    run_cli(repo, "approve-strategy", "--slug", slug)

    blocked_without_decision = run_cli_failure(repo, "generate-outline", "--slug", slug)
    assert blocked_without_decision.returncode != 0
    assert "decide-discovery" in (blocked_without_decision.stdout + blocked_without_decision.stderr)

    recorded_decision = parse_kv(run_cli(repo, "decide-discovery", "--slug", slug, "--decision", "no").stdout)
    assert recorded_decision["discovery_decision"] == "no"

    generated_outline = parse_kv(run_cli(repo, "generate-outline", "--slug", slug).stdout)
    assert generated_outline["outline"] == str(repo / "posts" / slug / "outline.md")

    resolve_generated_outline(repo / "posts" / slug / "outline.md")

    blocked_draft = run_cli_failure(repo, "generate-draft", "--slug", slug)
    assert blocked_draft.returncode != 0
    assert "approve-outline" in (blocked_draft.stdout + blocked_draft.stderr)

    run_cli(repo, "approve-outline", "--slug", slug)
    generated_draft = parse_kv(run_cli(repo, "generate-draft", "--slug", slug).stdout)
    assert generated_draft["draft"] == str(repo / "posts" / slug / "draft.qmd")


def test_generate_outline_requires_discovery_output_when_decision_is_yes(repo_root: Path, tmp_path: Path, monkeypatch) -> None:
    repo = copy_repo(repo_root, tmp_path)
    monkeypatch.chdir(repo)

    slug = "discovery-gated"
    run_cli(repo, "new-post", "--slug", slug, "--topic", "Discovery-gated post")
    strategy_path = repo / "posts" / slug / "strategy.md"
    resolve_generated_strategy(strategy_path)
    run_cli(repo, "approve-strategy", "--slug", slug)
    run_cli(repo, "decide-discovery", "--slug", slug, "--decision", "yes")

    blocked_outline = run_cli_failure(repo, "generate-outline", "--slug", slug)
    assert blocked_outline.returncode != 0
    assert f"posts/{slug}/inputs/discovery/discovery.md" in (blocked_outline.stdout + blocked_outline.stderr)

    discovery_summary = repo / "posts" / slug / "inputs" / "discovery" / "discovery.md"
    discovery_summary.parent.mkdir(parents=True, exist_ok=True)
    discovery_summary.write_text("# Discovery\n\nThe operator chose to run discovery.\n", encoding="utf-8")

    generated_outline = parse_kv(run_cli(repo, "generate-outline", "--slug", slug).stdout)
    assert generated_outline["outline"] == str(repo / "posts" / slug / "outline.md")


def test_generate_draft_quotes_yaml_sensitive_titles(repo_root: Path, tmp_path: Path, monkeypatch) -> None:
    repo = copy_repo(repo_root, tmp_path)
    monkeypatch.chdir(repo)

    slug = "colon-title"
    run_cli(repo, "new-post", "--slug", slug, "--topic", "Claude terms: connectors, skills, plugins, and MCP")
    strategy_path = repo / "posts" / slug / "strategy.md"
    resolve_generated_strategy(strategy_path)
    run_cli(repo, "approve-strategy", "--slug", slug)
    run_cli(repo, "decide-discovery", "--slug", slug, "--decision", "no")
    run_cli(repo, "generate-outline", "--slug", slug)
    resolve_generated_outline(repo / "posts" / slug / "outline.md")
    run_cli(repo, "approve-outline", "--slug", slug)
    run_cli(repo, "generate-draft", "--slug", slug)

    draft_path = repo / "posts" / slug / "draft.qmd"
    frontmatter, _ = extract_frontmatter(draft_path.read_text(encoding="utf-8"))
    assert frontmatter["title"] == "Claude terms: connectors, skills, plugins, and MCP"

    staged = parse_kv(run_cli(repo, "stage-attempt", "--slug", slug).stdout)
    attempt_root = repo / "posts" / slug / "runs" / staged["run_id"] / "attempts" / staged["attempt_id"]
    assert attempt_root.exists()


def test_prepare_visuals_writes_requests_without_editing_draft(repo_root: Path, tmp_path: Path, monkeypatch) -> None:
    repo = copy_repo(repo_root, tmp_path)
    monkeypatch.chdir(repo)

    slug = "visual-prep"
    build_post_with_draft(repo, slug)

    draft_path = repo / "posts" / slug / "draft.qmd"
    original_draft = draft_path.read_text(encoding="utf-8")
    draft_path.write_text(
        original_draft
        + "\n## Visual Section\n\n<!-- visual: show the pipeline stages -->\n\n## Second Visual\n\n<!-- visual: compare raw and prepared inputs -->\n",
        encoding="utf-8",
    )

    generated = parse_kv(run_cli(repo, "prepare-visuals", "--slug", slug).stdout)
    requests_path = repo / "posts" / slug / "visuals" / "requests.json"
    assert generated["requests"] == str(requests_path)
    assert generated["visual_count"] == "2"
    assert draft_path.read_text(encoding="utf-8").count("<!-- visual:") == 2

    payload = read_json(requests_path)
    assert payload["prepared_input"] == f"posts/{slug}/inputs/prepared/input.md"
    assert payload["brand_guide"] == "presets/default/brand_guide.md"
    assert payload["visual_identity"].startswith("## Visual Identity")
    assert "### Colour Tokens" in payload["visual_identity"]
    assert "### Aspect Ratios" in payload["visual_identity"]
    assert payload["visual_requirements"][0].startswith("Produce a self-contained HTML document")
    assert [item["verifier"] for item in payload["visual_verifiers"]["must_have"]] == [
        "visual_relevance",
        "text_visual_alignment",
        "source_attribution",
        "layout_integrity",
    ]
    assert payload["visual_verifiers"]["must_have"][0]["prompt_path"] == "prompts/visual_verifiers/visual_relevance.md"
    assert payload["visual_verifiers"]["must_have"][0]["rubric"].startswith("# Visual Relevance")
    assert len(payload["requests"]) == 2


def test_embed_visuals_replaces_selected_markers(repo_root: Path, tmp_path: Path, monkeypatch) -> None:
    repo = copy_repo(repo_root, tmp_path)
    monkeypatch.chdir(repo)

    slug = "visual-embed"
    build_post_with_draft(repo, slug)
    draft_path = repo / "posts" / slug / "draft.qmd"
    draft_path.write_text(
        draft_path.read_text(encoding="utf-8")
        + "\n## Visual Section\n\n<!-- visual: show the pipeline stages -->\n\n## Second Visual\n\n<!-- visual: compare raw and prepared inputs -->\n",
        encoding="utf-8",
    )

    run_cli(repo, "prepare-visuals", "--slug", slug)
    payload = read_json(repo / "posts" / slug / "visuals" / "requests.json")
    first_request = payload["requests"][0]
    first_html = repo / "posts" / slug / first_request["html_path"]
    first_html.parent.mkdir(parents=True, exist_ok=True)
    first_html.write_text("<html><body>visual 1</body></html>\n", encoding="utf-8")

    run_cli(repo, "embed-visuals", "--slug", slug, "--visual-id", first_request["visual_id"])
    rewritten = draft_path.read_text(encoding="utf-8")
    assert first_request["html_path"] in rewritten
    assert "<!-- visual: compare raw and prepared inputs -->" in rewritten


def test_approve_outline_rejects_generic_outline_headings(repo_root: Path, tmp_path: Path, monkeypatch) -> None:
    repo = copy_repo(repo_root, tmp_path)
    monkeypatch.chdir(repo)

    slug = "generic-outline"
    run_cli(repo, "new-post", "--slug", slug, "--topic", "Generic outline headings")
    strategy_path = repo / "posts" / slug / "strategy.md"
    resolve_generated_strategy(strategy_path)
    run_cli(repo, "approve-strategy", "--slug", slug)
    run_cli(repo, "decide-discovery", "--slug", slug, "--decision", "no")
    run_cli(repo, "generate-outline", "--slug", slug)

    outline_path = repo / "posts" / slug / "outline.md"
    frontmatter, _ = extract_frontmatter(outline_path.read_text(encoding="utf-8"))
    outline_path.write_text(
        format_markdown_with_frontmatter(
            frontmatter,
            "# Outline\n\n## Hook\n- Lead point.\n\n## Body section 1\n- Explain the distinction.\n",
        ),
        encoding="utf-8",
    )

    blocked = run_cli_failure(repo, "approve-outline", "--slug", slug)
    assert blocked.returncode != 0
    assert "publishable, reader-facing section title" in (blocked.stdout + blocked.stderr)
