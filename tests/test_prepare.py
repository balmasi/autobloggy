from __future__ import annotations

import os
import subprocess
import sys
from pathlib import Path

from pptx import Presentation

from autobloggy.artifacts import extract_frontmatter, read_sources
from autobloggy.prepare import run_prepare
from tests.helpers import copy_repo, resolve_generated_brief


def test_prepare_supports_markdown_seed(repo_root: Path, tmp_path: Path, monkeypatch) -> None:
    repo = copy_repo(repo_root, tmp_path)
    monkeypatch.chdir(repo)
    slug = "prepare-markdown"
    generated = run_prepare(slug, repo / "tests/fixtures/example_seed.md", "draft")

    assert "draft" in generated
    assert (repo / "posts" / slug / "claims.yaml").exists()


def test_prepare_supports_seed_directory_with_supporting_files(repo_root: Path, tmp_path: Path, monkeypatch) -> None:
    repo = copy_repo(repo_root, tmp_path)
    monkeypatch.chdir(repo)
    slug = "prepare-seed-dir"
    seed_root = repo / "posts" / slug / "seed"
    seed_root.mkdir(parents=True, exist_ok=True)
    seed_path = seed_root / "seed.md"
    seed_path.write_text(
        """---
title: Seed Directory Example
---

# Seed Directory Example

Seed notes can live beside supporting files instead of cluttering the post root.
""",
        encoding="utf-8",
    )
    (seed_root / "slides").mkdir()
    (seed_root / "slides" / "overview.txt").write_text("Slide notes.", encoding="utf-8")

    generated = run_prepare(slug, seed_root, "claims")

    assert "claims" in generated

    brief_frontmatter, _ = extract_frontmatter((repo / "posts" / slug / "brief.md").read_text(encoding="utf-8"))
    assert brief_frontmatter["seed_path"] == str(seed_path)
    assert brief_frontmatter["seed_root"] == str(seed_root)

    sources = read_sources(repo / "posts" / slug / "sources.yaml")
    seed_source = next(source for source in sources.sources if source.id == "src-seed")
    assert seed_source.locator == str(seed_path)
    assert str(seed_root) in (seed_source.notes or "")


def test_prepare_supports_pptx_seed(repo_root: Path, tmp_path: Path, monkeypatch) -> None:
    repo = copy_repo(repo_root, tmp_path)
    monkeypatch.chdir(repo)
    pptx_path = tmp_path / "seed.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Checklist rollout"
    slide.placeholders[1].text = "Handoffs dropped from 30 minutes to 5 minutes."
    presentation.save(pptx_path)

    slug = "prepare-pptx"
    generated = run_prepare(slug, pptx_path, "claims")

    assert "claims" in generated
    assert (repo / "posts" / slug / "sources.yaml").exists()


def test_prepare_cleans_markdown_seed_before_generating_artifacts(repo_root: Path, tmp_path: Path, monkeypatch) -> None:
    repo = copy_repo(repo_root, tmp_path)
    monkeypatch.chdir(repo)
    seed_path = repo / "tests" / "fixtures" / "noisy_seed.md"
    seed_path.write_text(
        """---
title: Noisy Seed
---

# Noisy Seed

Claude Code works in a terminal and can edit code across files.

NotebookLM works from a bounded source set and helps with synthesis before implementation. https://example.com/notebooklm

## Questions to answer

- What should move from research into execution?
- Where should implementation start?
""",
        encoding="utf-8",
    )

    slug = "prepare-noisy-markdown"
    run_prepare(slug, seed_path, "draft")

    brief_text = (repo / "posts" / slug / "brief.md").read_text(encoding="utf-8")
    outline_text = (repo / "posts" / slug / "outline.md").read_text(encoding="utf-8")
    claims_text = (repo / "posts" / slug / "claims.yaml").read_text(encoding="utf-8")
    draft_text = (repo / "posts" / slug / "draft.qmd").read_text(encoding="utf-8")
    sources_text = (repo / "posts" / slug / "sources.yaml").read_text(encoding="utf-8")

    assert "# Noisy Seed\n\n# Noisy Seed" not in draft_text
    assert draft_text.count("# Noisy Seed") == 1
    assert "## Noisy Seed" not in draft_text
    assert "Questions to answer" not in outline_text
    assert "https://example.com/notebooklm" not in claims_text
    assert sources_text.count("locator: https://example.com/notebooklm") == 1


def test_generate_brief_includes_required_voice_sections(repo_root: Path, tmp_path: Path, monkeypatch) -> None:
    repo = copy_repo(repo_root, tmp_path)
    monkeypatch.chdir(repo)
    slug = "prepare-brief-template"
    seed_path = repo / "tests" / "fixtures" / "example_seed.md"

    run_prepare(slug, seed_path, "brief")

    brief_text = (repo / "posts" / slug / "brief.md").read_text(encoding="utf-8")
    assert "## Target Voice" in brief_text
    assert "## Style Guardrails" in brief_text
    assert "## Evidence Standards" in brief_text
    assert "## Approval Checklist" in brief_text
    assert "[REQUIRED:" in brief_text
    assert "seed_root:" in brief_text


def test_cli_requires_brief_approval_before_generating_draft(repo_root: Path, tmp_path: Path) -> None:
    repo = copy_repo(repo_root, tmp_path)
    slug = "review-gated"
    seed_path = repo / "tests" / "fixtures" / "example_seed.md"

    first_prepare = subprocess.run(
        [sys.executable, "-m", "autobloggy.cli", "prepare", "--slug", slug, "--seed", str(seed_path)],
        cwd=repo,
        env={**os.environ, "PYTHONPATH": str(repo / "src")},
        text=True,
        capture_output=True,
        check=True,
    )
    assert "brief\t" in first_prepare.stdout
    assert not (repo / "posts" / slug / "draft.qmd").exists()

    brief_path = repo / "posts" / slug / "brief.md"
    brief_path.write_text(brief_path.read_text(encoding="utf-8") + "\nCustom review note.\n", encoding="utf-8")

    rejected_approval = subprocess.run(
        [sys.executable, "-m", "autobloggy.cli", "approve-brief", "--slug", slug],
        cwd=repo,
        env={**os.environ, "PYTHONPATH": str(repo / "src")},
        text=True,
        capture_output=True,
        check=False,
    )
    assert rejected_approval.returncode != 0
    assert "Brief is incomplete" in (rejected_approval.stdout + rejected_approval.stderr)

    resolve_generated_brief(brief_path)

    blocked_prepare = subprocess.run(
        [sys.executable, "-m", "autobloggy.cli", "prepare", "--slug", slug, "--seed", str(seed_path), "--through", "draft"],
        cwd=repo,
        env={**os.environ, "PYTHONPATH": str(repo / "src")},
        text=True,
        capture_output=True,
        check=False,
    )
    assert blocked_prepare.returncode != 0
    assert "approve-brief" in (blocked_prepare.stdout + blocked_prepare.stderr)

    subprocess.run(
        [sys.executable, "-m", "autobloggy.cli", "approve-brief", "--slug", slug],
        cwd=repo,
        env={**os.environ, "PYTHONPATH": str(repo / "src")},
        text=True,
        capture_output=True,
        check=True,
    )
    assert "status: approved" in brief_path.read_text(encoding="utf-8")

    final_prepare = subprocess.run(
        [sys.executable, "-m", "autobloggy.cli", "prepare", "--slug", slug, "--seed", str(seed_path), "--through", "draft"],
        cwd=repo,
        env={**os.environ, "PYTHONPATH": str(repo / "src")},
        text=True,
        capture_output=True,
        check=True,
    )
    assert "draft\t" in final_prepare.stdout
    assert (repo / "posts" / slug / "draft.qmd").exists()
    assert "Custom review note." in brief_path.read_text(encoding="utf-8")


def test_cli_accepts_seed_directory(repo_root: Path, tmp_path: Path) -> None:
    repo = copy_repo(repo_root, tmp_path)
    slug = "review-seed-dir"
    seed_root = repo / "posts" / slug / "seed"
    seed_root.mkdir(parents=True, exist_ok=True)
    seed_path = seed_root / "seed.md"
    seed_path.write_text(
        """# Seed Dir CLI

Putting the seed under posts/<slug>/seed keeps related files together.
""",
        encoding="utf-8",
    )
    (seed_root / "images").mkdir()
    (seed_root / "images" / "frame.txt").write_text("image placeholder", encoding="utf-8")

    first_prepare = subprocess.run(
        [sys.executable, "-m", "autobloggy.cli", "prepare", "--slug", slug, "--seed", str(seed_root)],
        cwd=repo,
        env={**os.environ, "PYTHONPATH": str(repo / "src")},
        text=True,
        capture_output=True,
        check=True,
    )
    assert "brief\t" in first_prepare.stdout

    brief_path = repo / "posts" / slug / "brief.md"
    assert f"seed_path: {seed_path}" in brief_path.read_text(encoding="utf-8")
