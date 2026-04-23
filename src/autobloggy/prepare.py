from __future__ import annotations

import re
import shutil
from pathlib import Path

from pptx import Presentation

from .artifacts import (
    extract_frontmatter,
    format_markdown_with_frontmatter,
    post_paths,
    read_text,
    read_yaml,
    text_fingerprint,
    write_text,
    write_yaml,
)
from .models import InputManifest, InputTextSource, InputVisualSource
from .presets import default_preset_name, preset_paths, presets_root, repo_relative_path, strategy_frontmatter_preset
from .utils import ensure_dir, now_iso, sentences, slugify


class InputContent:
    def __init__(self, title: str, text: str, headings: list[str], urls: list[str], kind: str):
        self.title = title
        self.text = text
        self.headings = headings
        self.urls = urls
        self.kind = kind


REQUIRED_STRATEGY_SECTIONS = (
    "Preset Context",
    "Core Question",
    "Audience",
    "Reader Outcome",
    "Target Voice",
    "Style Guardrails",
    "Must Cover",
    "Must Avoid",
    "Open Questions Before Approval",
    "Approval Checklist",
)

NON_PUBLISHABLE_OUTLINE_HEADINGS = {
    "hook",
    "opening",
    "context",
    "thesis",
    "problem",
    "what changed",
    "what it means",
    "decision table",
    "implications",
    "closing",
    "conclusion",
    "drafting notes",
    "notes for drafting",
}

UNRESOLVED_STRATEGY_MARKERS = (
    "[REQUIRED:",
    "- [ ]",
)

TEXT_INPUT_SUFFIXES = {".md", ".markdown", ".qmd", ".txt"}
PRESENTATION_SUFFIXES = {".pptx"}
SUPPORTED_INPUT_SUFFIXES = TEXT_INPUT_SUFFIXES | PRESENTATION_SUFFIXES
RAW_VISUAL_SUFFIXES = {".png", ".jpg", ".jpeg", ".svg", ".webp", ".gif"}
PREFERRED_INPUT_FILENAMES = (
    "input.md",
    "input.markdown",
    "input.qmd",
    "input.txt",
    "input.pptx",
)
INPUT_BUNDLE_EXCLUDED_HEADINGS = {
    "brief",
    "input sources",
    "raw text sources",
    "extracted text sources",
    "raw visual sources",
    "extracted visual sources",
    "text source notes",
    "supporting material notes",
    "links",
    "notes",
    "sources",
}


def markdown_headings(text: str) -> list[tuple[int, str]]:
    matches = re.finditer(r"^(#{1,6})\s+(.+)$", text, flags=re.MULTILINE)
    return [(len(match.group(1)), match.group(2).strip()) for match in matches]


def clean_markdown_input_text(text: str) -> str:
    cleaned = re.sub(r"<!--[\s\S]*?-->", "", text)
    lines = [line for line in cleaned.splitlines() if not re.match(r"^\s*#{1,6}\s+", line)]
    cleaned = "\n".join(lines)
    cleaned = re.sub(r"https?://[^\s)]+", "", cleaned)
    cleaned = re.sub(r"[ \t]+\n", "\n", cleaned)
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    return cleaned.strip()


def usable_markdown_headings(text: str, title: str) -> list[str]:
    headings: list[str] = []
    seen: set[str] = set()
    title_slug = slugify(title)
    for level, heading in markdown_headings(text):
        heading_slug = slugify(heading)
        if heading_slug == title_slug:
            continue
        if heading.casefold() in INPUT_BUNDLE_EXCLUDED_HEADINGS:
            continue
        if level == 1:
            continue
        if heading_slug in seen:
            continue
        seen.add(heading_slug)
        headings.append(heading)
    return headings


def humanize_name(value: str) -> str:
    cleaned = re.sub(r"[-_]+", " ", value).strip()
    return cleaned or "Untitled post"


def dedupe_preserving_order(values: list[str]) -> list[str]:
    seen: set[str] = set()
    deduped: list[str] = []
    for value in values:
        normalized = value.strip()
        if not normalized:
            continue
        key = normalized.casefold()
        if key in seen:
            continue
        seen.add(key)
        deduped.append(normalized)
    return deduped


def is_relative_to(path: Path, parent: Path) -> bool:
    try:
        path.resolve().relative_to(parent.resolve())
        return True
    except ValueError:
        return False


def post_relative_path(path: Path, post_root: Path) -> str:
    return str(path.resolve().relative_to(post_root.resolve()))


def supported_source_files(root: Path) -> list[Path]:
    if root.is_file():
        return [root] if root.suffix.lower() in SUPPORTED_INPUT_SUFFIXES else []
    if not root.exists():
        raise FileNotFoundError(f"Source path does not exist: {root}")
    return sorted(path for path in root.rglob("*") if path.is_file() and path.suffix.lower() in SUPPORTED_INPUT_SUFFIXES)


def resolve_input_path(input_candidate: Path) -> Path:
    if input_candidate.is_file():
        if input_candidate.suffix.lower() not in SUPPORTED_INPUT_SUFFIXES:
            raise ValueError(f"Unsupported input format: {input_candidate.suffix}")
        return input_candidate

    if not input_candidate.is_dir():
        raise FileNotFoundError(f"Input path does not exist: {input_candidate}")

    for filename in PREFERRED_INPUT_FILENAMES:
        candidate = input_candidate / filename
        if candidate.exists():
            return candidate

    candidates = sorted(
        path for path in input_candidate.iterdir() if path.is_file() and path.suffix.lower() in SUPPORTED_INPUT_SUFFIXES
    )
    if len(candidates) == 1:
        return candidates[0]
    if not candidates:
        raise ValueError(
            f"No supported main source file found in {input_candidate}. Add one of: {', '.join(PREFERRED_INPUT_FILENAMES)}."
        )
    raise ValueError(
        f"Multiple supported main source files found in {input_candidate}. Keep a single main source file at the top level."
    )


def parse_input(input_path: Path) -> InputContent:
    input_path = resolve_input_path(input_path)
    suffix = input_path.suffix.lower()
    if suffix in TEXT_INPUT_SUFFIXES:
        raw = read_text(input_path)
        frontmatter, body = extract_frontmatter(raw)
        body_headings = markdown_headings(body)
        title = str(frontmatter.get("title") or (body_headings[0][1] if body_headings else input_path.stem))
        urls = re.findall(r"https?://[^\s)]+", raw)
        cleaned_text = clean_markdown_input_text(body)
        headings = usable_markdown_headings(body, title)
        return InputContent(title=title, text=cleaned_text, headings=headings, urls=urls, kind="markdown")

    if suffix in PRESENTATION_SUFFIXES:
        presentation = Presentation(input_path)
        slide_titles: list[str] = []
        fragments: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            slide_texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                slide_titles.append(slide_texts[0].splitlines()[0].strip())
                fragments.append(f"Slide {index}: " + " | ".join(slide_texts))
        text = "\n".join(fragments)
        urls = re.findall(r"https?://[^\s)]+", text)
        title = slide_titles[0] if slide_titles else input_path.stem
        return InputContent(title=title, text=text, headings=slide_titles, urls=urls, kind="pptx")

    raise ValueError(f"Unsupported input format: {input_path.suffix}")


def summarize_text(text: str, max_sentences: int = 3) -> str:
    return " ".join(sentences(text)[:max_sentences]).strip()


def section_slug(heading: str) -> str:
    return slugify(heading) or "section"


def default_must_cover(input_content: InputContent) -> list[str]:
    headings = [heading for heading in input_content.headings if heading.strip()]
    if headings:
        return [heading.rstrip(".") for heading in headings[:5]]
    return [
        f"The real problem or decision behind {input_content.title}",
        "How the workflow, system, or argument works in practice",
        "The key tradeoffs, limits, and failure modes",
    ]


def required_strategy_section_issues(strategy_body: str) -> list[str]:
    issues: list[str] = []
    for heading in REQUIRED_STRATEGY_SECTIONS:
        if not re.search(rf"^##\s+{re.escape(heading)}\s*$", strategy_body, flags=re.MULTILINE):
            issues.append(f"Missing required section: {heading}.")
    return issues


def strategy_approval_issues(strategy_text: str) -> list[str]:
    _, body = extract_frontmatter(strategy_text)
    issues = required_strategy_section_issues(body)

    for marker in UNRESOLVED_STRATEGY_MARKERS:
        if marker in body:
            if marker == "[REQUIRED:":
                issues.append("Resolve every `[REQUIRED: ...]` prompt before approval.")
            elif marker == "- [ ]":
                issues.append("Check every item in the approval checklist before approval.")

    return issues


def outline_approval_issues(outline_text: str) -> list[str]:
    _, body = extract_frontmatter(outline_text)
    issues: list[str] = []
    headings = re.findall(r"^##\s+(.+)$", body, flags=re.MULTILINE)
    if not headings:
        issues.append("Outline has no sections. Add at least one ## heading before approving.")
        return issues

    flagged_headings: list[str] = []
    for heading in headings:
        normalized = heading.strip()
        lowered = normalized.casefold()
        if lowered in NON_PUBLISHABLE_OUTLINE_HEADINGS or lowered.startswith("body section "):
            flagged_headings.append(normalized)

    for heading in flagged_headings:
        issues.append(
            f"Replace generic outline heading `{heading}` with a publishable, reader-facing section title."
        )
    return issues


def extract_section_bullets(body: str, heading: str) -> list[str]:
    match = re.search(rf"^## {re.escape(heading)}\s*$([\s\S]*?)(?=^## |\Z)", body, flags=re.MULTILINE)
    if not match:
        return []
    return [item.strip() for item in re.findall(r"^\-\s+(.+)$", match.group(1), flags=re.MULTILINE)]


def render_strategy_template(template_text: str, context: dict[str, str]) -> str:
    rendered = template_text
    for key, value in context.items():
        rendered = rendered.replace(f"{{{{{key}}}}}", value)

    unresolved = sorted(set(re.findall(r"\{\{([a-z0-9_]+)\}\}", rendered)))
    if unresolved:
        raise ValueError("Strategy template has unresolved tokens: " + ", ".join(unresolved))
    return rendered.strip() + "\n"


def render_user_input_readme() -> str:
    return "\n".join(
        [
            "# User Inputs",
            "",
            "Use this folder only for human-owned material.",
            "",
            "- Put the plain-language post brief in `brief.md`.",
            "- Put original source files and folders under `raw/`.",
            "- Do not put generated files under `user_provided/`.",
            "- `../extracted/` contains deterministic extracts. It can be rebuilt.",
            "- `../prepared/` contains the canonical LLM-facing bundle. It can be rebuilt.",
            "",
        ]
    )


def render_brief_template() -> str:
    return "\n".join(
        [
            "<!--",
            "Add the plain-language brief here when you are not giving it conversationally.",
            "Include the topic, audience, must-cover points, and must-avoid framing.",
            "-->",
            "",
        ]
    )


def scaffold_input_layout(slug: str) -> dict[str, str]:
    paths = post_paths(slug)
    ensure_dir(paths.user_provided_root)
    ensure_dir(paths.user_raw_root)
    ensure_dir(paths.extracted_text_root)
    ensure_dir(paths.extracted_visual_root)
    ensure_dir(paths.prepared_root)
    ensure_dir(paths.discovery_root)

    if not paths.user_readme.exists():
        write_text(paths.user_readme, render_user_input_readme())
    if not paths.user_brief.exists():
        write_text(paths.user_brief, render_brief_template())

    return {
        "slug": slug,
        "brief": str(paths.user_brief),
        "raw": str(paths.user_raw_root),
        "prepared": str(paths.prepared_root),
    }


def render_brief_markdown(title: str | None, topic: str) -> str:
    heading = title or next((line.strip().rstrip(".") for line in topic.splitlines() if line.strip()), "Post brief")
    lines = [f"# {heading}", "", topic.strip(), ""]
    return "\n".join(lines)


def brief_has_substance(brief_path: Path) -> bool:
    if not brief_path.exists():
        return False
    raw = read_text(brief_path)
    _, body = extract_frontmatter(raw)
    cleaned = clean_markdown_input_text(body or raw)
    return bool(cleaned)


def load_brief_content(slug: str) -> InputContent | None:
    paths = post_paths(slug)
    if not brief_has_substance(paths.user_brief):
        return None
    return parse_input(paths.user_brief)


def load_input_manifest(manifest_path: Path) -> InputManifest:
    if not manifest_path.exists():
        return InputManifest(generated_at=now_iso())
    raw = read_yaml(manifest_path) or {}
    if not raw:
        return InputManifest(generated_at=now_iso())
    return InputManifest.model_validate(raw)


def write_input_manifest(manifest_path: Path, manifest: InputManifest) -> None:
    write_yaml(manifest_path, manifest.model_dump(mode="json"))


def legacy_user_files(slug: str) -> list[Path]:
    paths = post_paths(slug)
    if not paths.user_provided_root.exists():
        return []
    legacy_files: list[Path] = []
    for candidate in sorted(paths.user_provided_root.rglob("*")):
        if not candidate.is_file():
            continue
        relative = candidate.relative_to(paths.user_provided_root).as_posix()
        if relative in {"README.md", "brief.md", "input.md"}:
            continue
        if relative.startswith("raw/") or relative.startswith("supporting/"):
            continue
        legacy_files.append(candidate)
    return legacy_files


def raw_source_files(slug: str) -> list[Path]:
    paths = post_paths(slug)
    candidates: list[Path] = []
    for root in (paths.user_raw_root, paths.legacy_supporting_root):
        if root.exists():
            candidates.extend(path for path in root.rglob("*") if path.is_file())
    candidates.extend(legacy_user_files(slug))

    deduped: dict[Path, Path] = {}
    for candidate in candidates:
        deduped[candidate.resolve()] = candidate
    return sorted(deduped.values(), key=lambda path: post_relative_path(path, paths.root))


def enough_input_exists(slug: str) -> bool:
    paths = post_paths(slug)
    if brief_has_substance(paths.user_brief):
        return True
    if raw_source_files(slug):
        return True
    return paths.legacy_main_input.exists()


def candidate_title_from_content(content: InputContent | None) -> str | None:
    if content is None:
        return None
    normalized = content.title.strip().rstrip(".")
    if normalized and normalized.casefold() not in {"brief", "input"}:
        return normalized
    for line in content.text.splitlines():
        cleaned = line.strip().rstrip(".")
        if cleaned:
            return cleaned
    return None


def derived_title(
    brief_content: InputContent | None,
    raw_text_contents: list[tuple[InputTextSource, InputContent]],
    extracted_text_contents: list[tuple[InputTextSource, InputContent]],
    raw_visual_sources: list[InputVisualSource],
    fallback: str = "Untitled post",
) -> str:
    brief_title = candidate_title_from_content(brief_content)
    if brief_title:
        return brief_title
    for entry, content in [*raw_text_contents, *extracted_text_contents]:
        title = candidate_title_from_content(content)
        if title:
            return title
        if entry.title:
            return entry.title
    if raw_visual_sources:
        return humanize_name(Path(raw_visual_sources[0].path).stem)
    return fallback


def copy_user_source(source: Path, destination_root: Path) -> Path:
    source = source.resolve()
    destination_root = destination_root.resolve()

    if is_relative_to(source, destination_root):
        return source

    ensure_dir(destination_root)
    target = destination_root / source.name
    if source.is_dir():
        shutil.copytree(source, target, dirs_exist_ok=True)
        return target
    if source.is_file():
        shutil.copy2(source, target)
        return target
    raise FileNotFoundError(f"Source path does not exist: {source}")


def normalize_artifact_stem(path: Path, post_root: Path) -> str:
    relative = post_relative_path(path, post_root)
    return slugify(str(Path(relative).with_suffix("")))


def visual_source_id(relative_path: str, kind: str, source_file: str, source_locator: str | None) -> str:
    stem = slugify(Path(relative_path).stem)[:32]
    suffix = text_fingerprint(
        {
            "path": relative_path,
            "kind": kind,
            "source_file": source_file,
            "source_locator": source_locator or "",
        }
    )[:8]
    return f"{stem}-{suffix}"


def existing_visual_metadata(manifest: InputManifest) -> dict[str, InputVisualSource]:
    entries = [*manifest.raw_visual_sources, *manifest.extracted_visual_sources]
    return {entry.path: entry for entry in entries}


def build_visual_source(
    *,
    path: Path,
    kind: str,
    source_file: str,
    post_root: Path,
    source_locator: str | None = None,
    width_px: int | None = None,
    height_px: int | None = None,
    existing: InputVisualSource | None = None,
) -> InputVisualSource:
    relative_path = post_relative_path(path, post_root)
    return InputVisualSource(
        id=(existing.id if existing else visual_source_id(relative_path, kind, source_file, source_locator)),
        path=relative_path,
        kind=kind,
        source_file=source_file,
        source_locator=source_locator,
        width_px=width_px if width_px is not None else (existing.width_px if existing else None),
        height_px=height_px if height_px is not None else (existing.height_px if existing else None),
        caption=existing.caption if existing else "",
        description=existing.description if existing else "",
        tags=list(existing.tags) if existing else [],
    )


def write_extracted_text(path: Path, content: InputContent, slug: str) -> Path:
    paths = post_paths(slug)
    output_path = paths.extracted_text_root / f"{normalize_artifact_stem(path, paths.root)}.md"
    body_lines = [f"# {content.title}", ""]
    if content.text.strip():
        body_lines.extend([content.text.strip(), ""])
    else:
        body_lines.extend([f"Deterministic extract from `{post_relative_path(path, paths.root)}`.", ""])
    frontmatter = {
        "title": content.title,
        "source_file": post_relative_path(path, paths.root),
        "generated_at": now_iso(),
    }
    write_text(output_path, format_markdown_with_frontmatter(frontmatter, "\n".join(body_lines).strip() + "\n"))
    return output_path


def extract_pptx_visuals(path: Path, slug: str, existing_by_path: dict[str, InputVisualSource]) -> list[InputVisualSource]:
    paths = post_paths(slug)
    presentation = Presentation(path)
    output_root = ensure_dir(paths.extracted_visual_root / normalize_artifact_stem(path, paths.root))
    source_file = post_relative_path(path, paths.root)
    visual_sources: list[InputVisualSource] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        picture_index = 0
        for shape in slide.shapes:
            if not hasattr(shape, "image"):
                continue
            picture_index += 1
            image = shape.image
            extension = image.ext or "bin"
            output_path = output_root / f"slide-{slide_index:02d}-{picture_index:02d}.{extension}"
            output_path.write_bytes(image.blob)
            width_px: int | None = None
            height_px: int | None = None
            try:
                width_px, height_px = image.size
            except Exception:
                width_px, height_px = None, None
            relative_path = post_relative_path(output_path, paths.root)
            visual_sources.append(
                build_visual_source(
                    path=output_path,
                    kind="pptx_extract",
                    source_file=source_file,
                    post_root=paths.root,
                    source_locator=f"slide {slide_index}, picture {picture_index}",
                    width_px=width_px,
                    height_px=height_px,
                    existing=existing_by_path.get(relative_path),
                )
            )

    return visual_sources


def render_prepared_input_markdown(
    *,
    title: str,
    slug: str,
    brief_content: InputContent | None,
    raw_text_contents: list[tuple[InputTextSource, InputContent]],
    extracted_text_contents: list[tuple[InputTextSource, InputContent]],
    raw_visual_sources: list[InputVisualSource],
    extracted_visual_sources: list[InputVisualSource],
) -> str:
    lines = [f"# {title}", ""]

    if brief_content and brief_content.text.strip():
        lines.extend(["## Brief", "", brief_content.text.strip(), ""])
    else:
        lines.extend(["This post starts from the files and extracts currently recorded for the post.", ""])

    lines.extend(["## Input Sources", ""])

    lines.extend(["### Raw Text Sources", ""])
    if raw_text_contents:
        lines.extend(f"- `{entry.path}`" for entry, _ in raw_text_contents)
    else:
        lines.append("- None")
    lines.append("")

    lines.extend(["### Extracted Text Sources", ""])
    if extracted_text_contents:
        lines.extend(f"- `{entry.path}` (from `{entry.extracted_from}`)" for entry, _ in extracted_text_contents)
    else:
        lines.append("- None")
    lines.append("")

    lines.extend(["### Raw Visual Sources", ""])
    if raw_visual_sources:
        lines.extend(f"- `{entry.path}`" for entry in raw_visual_sources)
    else:
        lines.append("- None")
    lines.append("")

    lines.extend(["### Extracted Visual Sources", ""])
    if extracted_visual_sources:
        lines.extend(f"- `{entry.path}` (from `{entry.source_file}`)" for entry in extracted_visual_sources)
    else:
        lines.append("- None")
    lines.append("")

    all_text_sources = [*raw_text_contents, *extracted_text_contents]
    if all_text_sources:
        lines.extend(["## Text Source Notes", ""])
        for entry, content in all_text_sources:
            label = content.title.strip() or humanize_name(Path(entry.path).stem)
            summary = summarize_text(content.text, 2) or f"Supporting material from `{entry.path}`."
            lines.extend([f"### {label}", "", f"Source: `{entry.path}`", "", summary, ""])

    frontmatter = {
        "title": title,
        "generated_at": now_iso(),
        "slug": slug,
    }
    return format_markdown_with_frontmatter(frontmatter, "\n".join(lines).strip() + "\n")


def manifest_text_source_paths(manifest: InputManifest, slug: str) -> list[Path]:
    paths = post_paths(slug)
    text_entries = [*manifest.raw_text_sources, *manifest.extracted_text_sources]
    resolved: list[Path] = []
    for entry in text_entries:
        resolved.append(paths.root / entry.path)
    return resolved


def prepare_post_inputs(slug: str, require_sources: bool = True) -> dict[str, str]:
    paths = post_paths(slug)
    scaffold_input_layout(slug)

    existing_manifest = load_input_manifest(paths.input_manifest)
    existing_visuals = existing_visual_metadata(existing_manifest)

    if paths.extracted_text_root.exists():
        shutil.rmtree(paths.extracted_text_root)
    if paths.extracted_visual_root.exists():
        shutil.rmtree(paths.extracted_visual_root)
    ensure_dir(paths.extracted_text_root)
    ensure_dir(paths.extracted_visual_root)
    ensure_dir(paths.prepared_root)

    brief_content = load_brief_content(slug)
    discovered_raw_files = raw_source_files(slug)

    raw_text_contents: list[tuple[InputTextSource, InputContent]] = []
    extracted_text_contents: list[tuple[InputTextSource, InputContent]] = []
    raw_visual_sources: list[InputVisualSource] = []
    extracted_visual_sources: list[InputVisualSource] = []

    for path in discovered_raw_files:
        suffix = path.suffix.lower()
        relative_path = post_relative_path(path, paths.root)
        if suffix in TEXT_INPUT_SUFFIXES:
            content = parse_input(path)
            raw_text_contents.append(
                (
                    InputTextSource(
                        path=relative_path,
                        kind=content.kind,
                        title=content.title,
                        headings=content.headings,
                    ),
                    content,
                )
            )
            continue

        if suffix in PRESENTATION_SUFFIXES:
            content = parse_input(path)
            extracted_path = write_extracted_text(path, content, slug)
            extracted_content = parse_input(extracted_path)
            extracted_text_contents.append(
                (
                    InputTextSource(
                        path=post_relative_path(extracted_path, paths.root),
                        kind="pptx_extract",
                        title=extracted_content.title,
                        headings=extracted_content.headings,
                        extracted_from=relative_path,
                    ),
                    extracted_content,
                )
            )
            extracted_visual_sources.extend(extract_pptx_visuals(path, slug, existing_visuals))
            continue

        if suffix in RAW_VISUAL_SUFFIXES:
            raw_visual_sources.append(
                build_visual_source(
                    path=path,
                    kind="raw_visual",
                    source_file=relative_path,
                    post_root=paths.root,
                    existing=existing_visuals.get(relative_path),
                )
            )

    if not brief_content and not raw_text_contents and not extracted_text_contents and not raw_visual_sources and not extracted_visual_sources:
        if paths.legacy_main_input.exists():
            legacy_content = parse_input(paths.legacy_main_input)
            raw_text_contents.append(
                (
                    InputTextSource(
                        path=post_relative_path(paths.legacy_main_input, paths.root),
                        kind="legacy_input_bundle",
                        title=legacy_content.title,
                        headings=legacy_content.headings,
                    ),
                    legacy_content,
                )
            )
        elif require_sources:
            raise ValueError(
                f"No post brief or raw source files found for `{slug}`. Add content to {paths.user_brief} or place files under {paths.user_raw_root}, then run `autobloggy prepare-inputs --slug {slug}`."
            )

    title = derived_title(brief_content, raw_text_contents, extracted_text_contents, raw_visual_sources, humanize_name(slug))
    prepared_input = render_prepared_input_markdown(
        title=title,
        slug=slug,
        brief_content=brief_content,
        raw_text_contents=raw_text_contents,
        extracted_text_contents=extracted_text_contents,
        raw_visual_sources=raw_visual_sources,
        extracted_visual_sources=extracted_visual_sources,
    )
    write_text(paths.prepared_input, prepared_input)

    manifest = InputManifest(
        generated_at=now_iso(),
        brief=(post_relative_path(paths.user_brief, paths.root) if brief_content else None),
        raw_text_sources=[entry for entry, _ in raw_text_contents],
        extracted_text_sources=[entry for entry, _ in extracted_text_contents],
        raw_visual_sources=raw_visual_sources,
        extracted_visual_sources=extracted_visual_sources,
        canonical_input=post_relative_path(paths.prepared_input, paths.root),
    )
    write_input_manifest(paths.input_manifest, manifest)

    return {
        "brief": str(paths.user_brief),
        "manifest": str(paths.input_manifest),
        "input": str(paths.prepared_input),
    }


def load_post_input_bundle(slug: str) -> tuple[Path, InputContent, list[Path]]:
    paths = post_paths(slug)
    if not paths.prepared_input.exists() or not paths.input_manifest.exists():
        prepare_post_inputs(slug, require_sources=True)
    input_content = parse_input(paths.prepared_input)
    manifest = load_input_manifest(paths.input_manifest)
    return paths.prepared_input, input_content, manifest_text_source_paths(manifest, slug)


def generate_strategy(
    slug: str,
    input_content: InputContent,
    input_path: Path,
    input_root: Path,
    preset_name: str,
    repo_context: Path,
) -> str:
    preset = preset_paths(preset_name, repo_context)
    must_cover = default_must_cover(input_content)
    strategy_template_path = repo_relative_path(preset.strategy_template, repo_context)
    writing_guide_path = repo_relative_path(preset.writing_guide, repo_context)
    brand_guide_path = repo_relative_path(preset.brand_guide, repo_context)
    frontmatter = {
        "slug": slug,
        "title": input_content.title,
        "input_path": str(input_path),
        "input_root": str(input_root),
        "input_type": input_content.kind,
        "input_manifest": repo_relative_path(post_paths(slug, repo_context).input_manifest, repo_context),
        "preset": preset.name,
        "preset_dir": repo_relative_path(preset.root, repo_context),
        "preset_strategy_template": strategy_template_path,
        "preset_writing_guide": writing_guide_path,
        "preset_brand_guide": brand_guide_path,
        "generated_at": now_iso(),
        "status": "needs_review",
    }
    body = render_strategy_template(
        read_text(preset.strategy_template),
        {
            "topic_summary": summarize_text(input_content.text) or f"{input_content.title} is the working topic for this post.",
            "title": input_content.title,
            "preset_name": preset.name,
            "strategy_template_path": strategy_template_path,
            "writing_guide_path": writing_guide_path,
            "brand_guide_path": brand_guide_path,
            "must_cover_bullets": "\n".join(f"- {item}" for item in must_cover),
        },
    )
    template_issues = required_strategy_section_issues(body)
    if template_issues:
        raise ValueError("Preset strategy template is incomplete:\n- " + "\n- ".join(template_issues))
    return format_markdown_with_frontmatter(frontmatter, body)


def generate_outline(strategy_text: str, input_content: InputContent) -> str:
    strategy_frontmatter, _ = extract_frontmatter(strategy_text)
    slug = strategy_frontmatter.get("slug", "")
    title = strategy_frontmatter.get("title", input_content.title)
    preset_name = str(strategy_frontmatter.get("preset") or default_preset_name())
    preset_dir = str(strategy_frontmatter.get("preset_dir") or repo_relative_path(preset_paths(preset_name).root))
    outline_frontmatter = {
        "slug": slug,
        "title": title,
        "preset": preset_name,
        "preset_dir": preset_dir,
        "generated_at": now_iso(),
        "status": "needs_review",
    }
    stub_body = f"# Outline\n\nTitle: {title}\n"
    return format_markdown_with_frontmatter(outline_frontmatter, stub_body)


def generate_draft(strategy_text: str, outline_text: str, input_content: InputContent) -> str:
    frontmatter, strategy_body = extract_frontmatter(strategy_text)
    title = frontmatter.get("title", "Untitled post")
    preset_name = str(frontmatter.get("preset") or default_preset_name())
    preset_dir = str(frontmatter.get("preset_dir") or repo_relative_path(preset_paths(preset_name).root))
    outline_headings = re.findall(r"^##\s+(.+)$", outline_text, flags=re.MULTILINE)
    source_sentences = [
        sentence.rstrip(".")
        for sentence in sentences(input_content.text)
        if len(sentence.split()) >= 6 and not sentence.endswith("?")
    ]
    if not source_sentences:
        source_sentences = [summarize_text(input_content.text, 1) or f"{title} is the focus of this draft."]
    cursor = 0

    draft_frontmatter = {
        "title": title,
        "format": "gfm",
        "preset": preset_name,
        "preset_dir": preset_dir,
    }

    lines = [
        f"# {title}",
        "",
    ]
    intro = summarize_text(strategy_body, 2) or f"{title} is the focus of this draft."
    lines.extend([intro, ""])

    for heading in outline_headings:
        if heading.lower() == "opening":
            continue
        lines.extend([f"## {heading}", ""])
        section_lines = source_sentences[cursor : cursor + 2]
        if not section_lines:
            section_lines = [f"Explain {heading.lower()} with concrete examples, tradeoffs, and operator detail."]
        cursor += len(section_lines)
        for section_line in section_lines:
            lines.append(f"{section_line}.")
            lines.append("")

    if "Conclusion" not in outline_headings:
        lines.extend(["## Conclusion", "", "Close with the practical takeaway.", ""])
    else:
        lines.append("Close with the practical takeaway.")
        lines.append("")

    return format_markdown_with_frontmatter(draft_frontmatter, "\n".join(lines).strip() + "\n")


def scaffold_preset(name: str, template_name: str = "default") -> Path:
    target_root = presets_root() / name
    if target_root.exists():
        raise ValueError(f"Preset `{name}` already exists at {repo_relative_path(target_root)}.")

    template = preset_paths(template_name)
    ensure_dir(target_root)
    for source_file in (template.strategy_template, template.writing_guide, template.brand_guide):
        shutil.copy2(source_file, target_root / source_file.name)
    return target_root


def run_new_post(
    slug: str | None,
    title: str | None,
    topic: str | None,
    source_paths: list[Path] | None = None,
    preset_name: str | None = None,
) -> dict[str, str]:
    source_paths = [path.resolve() for path in (source_paths or [])]
    if not slug and not title and not topic and not source_paths:
        raise ValueError("new-post requires a slug, a topic, a title, or at least one --source.")

    derived_slug = slug
    if not derived_slug:
        preview_title = title or next((path.stem for path in source_paths if path.name), None) or topic or "post"
        derived_slug = slugify(preview_title)

    paths = post_paths(derived_slug)
    if paths.strategy.exists() or paths.outline.exists() or paths.draft.exists():
        raise ValueError(f"Post `{derived_slug}` already has generated artifacts. Use the existing post commands for that post.")

    payload = scaffold_input_layout(derived_slug)
    if topic is not None or title is not None:
        write_text(paths.user_brief, render_brief_markdown(title, topic or ""))

    for source_path in source_paths:
        copy_user_source(source_path, paths.user_raw_root)

    if not enough_input_exists(derived_slug):
        payload["status"] = "scaffolded"
        payload["next_step"] = (
            f"Add a brief to {paths.user_brief} or place source files under {paths.user_raw_root}, "
            f"then run `autobloggy prepare-inputs --slug {derived_slug}`."
        )
        return payload

    prepared = prepare_post_inputs(derived_slug, require_sources=True)
    input_path, input_content, _ = load_post_input_bundle(derived_slug)
    resolved_preset_name = preset_name or default_preset_name(paths.root)
    strategy_text = generate_strategy(derived_slug, input_content, input_path, paths.prepared_root, resolved_preset_name, paths.root)
    write_text(paths.strategy, strategy_text)

    return {
        "slug": derived_slug,
        "preset": resolved_preset_name,
        "brief": str(paths.user_brief),
        "manifest": prepared["manifest"],
        "input": prepared["input"],
        "strategy": str(paths.strategy),
        "status": "generated",
    }


def run_generate_outline(slug: str) -> dict[str, str]:
    paths = post_paths(slug)
    if not paths.strategy.exists():
        raise FileNotFoundError(f"Strategy does not exist: {paths.strategy}")

    input_path, input_content, _ = load_post_input_bundle(slug)
    strategy_text = read_text(paths.strategy)
    frontmatter, _ = extract_frontmatter(strategy_text)
    if frontmatter.get("status") != "approved":
        raise ValueError(
            f"Strategy is not approved. Review posts/{slug}/strategy.md and run `autobloggy approve-strategy --slug {slug}` first."
        )
    discovery_decision = str(frontmatter.get("discovery_decision", "")).strip().lower()
    if discovery_decision not in {"yes", "no"}:
        raise ValueError(
            "Discovery decision is missing. Ask the user for an explicit yes/no decision and run "
            f"`autobloggy decide-discovery --slug {slug} --decision yes|no` before generating the outline."
        )
    if discovery_decision == "yes" and not paths.discovery_summary.exists():
        raise ValueError(
            f"Discovery was approved but not completed. Write discovery output to posts/{slug}/inputs/discovery/discovery.md before generating the outline."
        )

    outline_text = generate_outline(strategy_text, input_content)
    write_text(paths.outline, outline_text)
    return {
        "input": str(input_path),
        "outline": str(paths.outline),
    }


def run_generate_draft(slug: str) -> dict[str, str]:
    paths = post_paths(slug)
    if not paths.strategy.exists():
        raise FileNotFoundError(f"Strategy does not exist: {paths.strategy}")
    if not paths.outline.exists():
        raise FileNotFoundError(f"Outline does not exist: {paths.outline}")

    strategy_text = read_text(paths.strategy)
    strategy_frontmatter, _ = extract_frontmatter(strategy_text)
    if strategy_frontmatter.get("status") != "approved":
        raise ValueError(
            f"Strategy is not approved. Review posts/{slug}/strategy.md and run `autobloggy approve-strategy --slug {slug}` first."
        )

    outline_text = read_text(paths.outline)
    outline_frontmatter, _ = extract_frontmatter(outline_text)
    if outline_frontmatter.get("status") != "approved":
        raise ValueError(
            f"Outline is not approved. Review posts/{slug}/outline.md and run `autobloggy approve-outline --slug {slug}` first."
        )
    outline_issues = outline_approval_issues(outline_text)
    if outline_issues:
        raise ValueError("Outline is not draftable:\n- " + "\n- ".join(outline_issues))

    _, input_content, _ = load_post_input_bundle(slug)
    draft_text = generate_draft(strategy_text, outline_text, input_content)
    write_text(paths.draft, draft_text)
    return {"draft": str(paths.draft)}


def resolve_post_preset(strategy_path: Path, requested_preset: str | None) -> str:
    recorded_preset = strategy_frontmatter_preset(strategy_path)
    if requested_preset and recorded_preset and requested_preset != recorded_preset:
        raise ValueError(
            "Preset mismatch: strategy.md is using preset "
            f"`{recorded_preset}`, but this post was asked to use `{requested_preset}`. "
            "Create a new post with `autobloggy new-post --preset "
            f"{requested_preset}` or keep the recorded preset."
        )

    return requested_preset or recorded_preset or default_preset_name(strategy_path)
